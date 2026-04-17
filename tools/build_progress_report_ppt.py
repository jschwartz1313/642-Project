from __future__ import annotations

import sys
from pathlib import Path

sys.path.append("/Users/jakeschwartz/.local/lib/python3.13/site-packages")

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


TEMPLATE = Path("/Users/jakeschwartz/Downloads/642 Project Progress Report.pptx")
OUTPUT = Path("/Users/jakeschwartz/Downloads/642 Project Progress Report - Revised.pptx")


def set_text(shape, text: str) -> None:
    shape.text = text


def main() -> None:
    prs = Presentation(TEMPLATE)

    # Slide 1
    set_text(
        prs.slides[0].shapes[0],
        "Renewable Electricity Growth:\nRegression Analysis Progress Report",
    )
    set_text(prs.slides[0].shapes[1], "Jake Schwartz and Daniel Wood")

    # Slide 2
    slide = prs.slides[1]
    set_text(slide.shapes[3], "Project Overview")
    set_text(slide.shapes[4].shapes[1], "Primary Objective")
    set_text(
        slide.shapes[4].shapes[2],
        "Evaluate how wind and solar capacity relate to renewable electricity generation and renewable electricity share, then extend the project into forecasting.",
    )
    set_text(slide.shapes[5].shapes[1], "Methodology")
    set_text(
        slide.shapes[5].shapes[2],
        "Merged and cleaned 17 source files.\nUsed exploratory plots plus country and year fixed-effects regressions.\nForecasting is planned, but not yet implemented.",
    )
    set_text(slide.shapes[6].shapes[1], "Current Data Scope")
    set_text(
        slide.shapes[6].shapes[2],
        "Broad renewables panel: 251 entities.\nCapacity panel: 12 entities.\nFinal regression sample: 9 countries from 1997-2021.",
    )

    # Slide 3
    slide = prs.slides[2]
    set_text(slide.shapes[3], "Data Preparation")
    set_text(slide.shapes[5], "Processing Pipeline")
    set_text(
        slide.shapes[6],
        "Merged 17 CSV files into analytic panels, standardized variable names, and removed aggregate regions from the regression stage to keep the sample country-based.",
    )
    set_text(slide.shapes[8], "Broad Panel")
    set_text(slide.shapes[9], "251")
    set_text(slide.shapes[10], "Countries/regions with renewable generation and electricity-share measures.")
    set_text(slide.shapes[12], "Capacity Panel")
    set_text(slide.shapes[13], "12")
    set_text(slide.shapes[14], "Entities with installed wind and solar capacity variables before filtering.")
    set_text(slide.shapes[16], "Regression Sample")
    set_text(slide.shapes[17], "9")
    set_text(slide.shapes[18], "Countries used in the fixed-effects models after excluding 3 aggregate regions.")

    # Slide 4
    slide = prs.slides[3]
    set_text(slide.shapes[2], "Exploratory Findings")
    set_text(
        slide.shapes[3].shapes[1],
        "Visualization Methods\nTime-series plots track wind/solar capacity and generation by country.\nScatterplots pool capacity versus generation across all observations.",
    )
    set_text(
        slide.shapes[4].shapes[1],
        "Cross-Country Scale\nChina and the United States dominate absolute capacity and generation levels.\nThat scale gap supports using country fixed effects in the regressions.",
    )
    set_text(
        slide.shapes[5].shapes[1],
        "Temporal Trends\nWind capacity appears earlier in the sample.\nSolar adoption accelerates later, especially after about 2008.\nGrowth paths differ substantially across countries.",
    )

    # Slide 5
    slide = prs.slides[4]
    set_text(slide.shapes[2], "Regression Results")
    set_text(
        slide.shapes[3].shapes[1],
        "Regression 1: Capacity -> Generation\nWind generation model: 1 additional GW predicts about 1.96 TWh more wind generation (R² = 0.970, N = 225).\nSolar generation model: 1 additional GW predicts about 1.06 TWh more solar generation (R² = 0.973, N = 225).\nBoth coefficients are statistically significant.",
    )
    set_text(
        slide.shapes[4].shapes[1],
        "Regression 2: Capacity -> Renewable Share\nIn the joint model, wind and solar capacity coefficients are near zero and not statistically significant.\nModel fit remains moderate (R² = 0.651, N = 225), but within-country capacity growth alone does not explain short-run changes in renewable electricity share in this sample.",
    )

    # Slide 6
    slide = prs.slides[5]
    set_text(slide.shapes[2], "Challenges")
    set_text(
        slide.shapes[3].shapes[1],
        "Data Constraints\nThe fixed-effects regressions use only 9 countries from 1997-2021.\nThat keeps the sample consistent, but limits how broadly the results can be generalized.",
    )
    set_text(
        slide.shapes[4].shapes[1],
        "Interpretation Risk\nCapacity and generation are mechanically linked, so Regression 1 mainly validates the physical relationship between installed capacity and output.\nIt is less informative as a causal policy result.",
    )
    set_text(
        slide.shapes[5].shapes[1],
        "Specification Gaps\nNo lagged capacity terms yet.\nNo controls for demand, policy, storage, or grid constraints.\nForecasting work still needs to be built and validated.",
    )

    # Slide 7
    slide = prs.slides[6]
    set_text(slide.shapes[2], "Next Steps")
    set_text(
        slide.shapes[3].shapes[1],
        "Refine the models\nAdd 1-3 year lagged capacity terms and rerun the fixed-effects specifications.",
    )
    set_text(
        slide.shapes[4].shapes[1],
        "Strengthen interpretation\nAdd controls such as hydro share, electricity demand, or policy proxies if data are available.",
    )
    set_text(
        slide.shapes[5].shapes[1],
        "Stress-test the sample\nRun country-level sensitivity checks and compare joint versus single-technology specifications.",
    )
    set_text(
        slide.shapes[6].shapes[1],
        "Start forecasting\nDefine the target series, split training and test periods, and evaluate forecast accuracy out of sample.",
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
