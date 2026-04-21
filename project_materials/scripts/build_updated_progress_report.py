from __future__ import annotations

from pathlib import Path

from pptx import Presentation


REPO_ROOT = Path(__file__).resolve().parents[2]
MATERIALS_ROOT = Path(__file__).resolve().parents[1]

TEMPLATE = MATERIALS_ROOT / "presentations" / "642 Project Progress Report - Updated.pptx"
OUTPUT = MATERIALS_ROOT / "presentations" / "642 Project Progress Report - Updated.generated.pptx"

ARIMA_IMG = str(
    REPO_ROOT / "side_analyses" / "renewables_forecasting_clustering" / "outputs" / "arima_forecast.png"
)
KMEANS_IMG = str(
    REPO_ROOT / "side_analyses" / "renewables_forecasting_clustering" / "outputs" / "kmeans_clusters_pca.png"
)
DECOMP_IMG = str(
    REPO_ROOT / "focused_renewables_analysis" / "outputs" / "renewable_share_decomposition.png"
)
LAG_SHARE_IMG = str(
    REPO_ROOT / "focused_renewables_analysis" / "outputs" / "lagged_capacity_renew_share_elec.png"
)
LAG_WIND_IMG = str(
    REPO_ROOT / "focused_renewables_analysis" / "outputs" / "lagged_capacity_wind_gen.png"
)
LAG_SOLAR_IMG = str(
    REPO_ROOT / "focused_renewables_analysis" / "outputs" / "lagged_capacity_solar_gen.png"
)


def set_text(shape, text: str) -> None:
    shape.text = text


def remove_shape(shape) -> None:
    sp = shape._element
    sp.getparent().remove(sp)


def replace_picture(slide, old_shape, image_path: str) -> None:
    left, top, width, height = old_shape.left, old_shape.top, old_shape.width, old_shape.height
    remove_shape(old_shape)
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)


def main() -> None:
    prs = Presentation(TEMPLATE)

    # Slide 1
    slide = prs.slides[0]
    set_text(
        slide.shapes[0],
        "Renewable Electricity Growth:\nRegression, Forecasting, and Transition Analysis",
    )
    set_text(slide.shapes[1], "Jake Schwartz and Daniel Wood")

    # Slide 2
    slide = prs.slides[1]
    set_text(slide.shapes[3], "Project Overview")
    set_text(slide.shapes[4].shapes[1], "Current Objective")
    set_text(
        slide.shapes[4].shapes[2],
        "Explain renewable electricity growth with panel regressions, then extend the project with forecasting, clustering, decomposition, and lagged-capacity analysis.",
    )
    set_text(slide.shapes[5].shapes[1], "Methods Now Completed")
    set_text(
        slide.shapes[5].shapes[2],
        "Data merging and cleaning.\nExploratory analysis and fixed-effects regression.\nARIMA forecasting, K-means clustering, renewable-share decomposition, and lagged models.",
    )
    set_text(slide.shapes[6].shapes[1], "Core Data")
    set_text(
        slide.shapes[6].shapes[2],
        "17 renewable-energy source files merged into analytic panels.\nBroad panel for cross-country patterns and a smaller capacity panel for deeper modeling.",
    )

    # Slide 3
    slide = prs.slides[2]
    set_text(slide.shapes[3], "Data Preparation and Scope")
    set_text(slide.shapes[5], "Processing Pipeline")
    set_text(
        slide.shapes[6],
        "Merged 17 source CSVs, standardized variable names, checked country overlap, and removed aggregate regions from the modeling stage to keep the panel country-based.",
    )
    set_text(slide.shapes[8], "Broad Panel")
    set_text(slide.shapes[9], "251")
    set_text(slide.shapes[10], "Entities with renewable generation and electricity-share measures.")
    set_text(slide.shapes[12], "Capacity Panel")
    set_text(slide.shapes[13], "12")
    set_text(slide.shapes[14], "Entities with installed wind and solar capacity before filtering.")
    set_text(slide.shapes[16], "Regression Sample")
    set_text(slide.shapes[17], "9")
    set_text(slide.shapes[18], "Countries used in fixed-effects, lagged-capacity, and forecasting extensions.")

    # Slide 4
    slide = prs.slides[3]
    set_text(slide.shapes[3], "ARIMA Forecasting")
    set_text(slide.shapes[4].shapes[1], "Forecasting Result")
    set_text(
        slide.shapes[4].shapes[2],
        "ARIMA forecasting is now implemented for renewable electricity share. For Denmark, the selected model was ARIMA(0,2,2) with an 8-year holdout and MAPE of about 6.57%.",
    )
    replace_picture(slide, slide.shapes[5], ARIMA_IMG)

    # Slide 5
    slide = prs.slides[4]
    set_text(slide.shapes[3], "K-means Country Clusters")
    replace_picture(slide, slide.shapes[4], KMEANS_IMG)

    # Slide 6
    slide = prs.slides[5]
    set_text(slide.shapes[4], "Renewable Share Decomposition")
    replace_picture(slide, slide.shapes[3], DECOMP_IMG)

    # Slide 7
    slide = prs.slides[6]
    set_text(slide.shapes[3], "Lagged Capacity and Renewable Share")
    replace_picture(slide, slide.shapes[4], LAG_SHARE_IMG)

    # Slide 8
    slide = prs.slides[7]
    set_text(slide.shapes[2], "What EDA and Regression Already Showed")
    set_text(
        slide.shapes[3].shapes[1],
        "Exploratory Patterns\nWind generally appears earlier in the sample.\nSolar accelerates later, especially after about 2008.\nChina and the United States dominate absolute scale, which motivated fixed effects.",
    )
    set_text(
        slide.shapes[4].shapes[1],
        "Baseline Regression Results\nWind generation rises by about 1.96 TWh per added GW of wind capacity.\nSolar generation rises by about 1.06 TWh per added GW of solar capacity.\nBoth models fit very strongly with R-squared around 0.97.",
    )
    set_text(
        slide.shapes[5].shapes[1],
        "Interpretation\nCapacity clearly predicts later generation, but the joint renewable-share model remains weaker.\nThat is why the new extensions focus on timing, country profiles, and source decomposition.",
    )

    # Slide 9
    slide = prs.slides[8]
    set_text(slide.shapes[2], "Lagged Wind Generation Effects")
    replace_picture(slide, slide.shapes[3], LAG_WIND_IMG)

    # Slide 10
    slide = prs.slides[9]
    set_text(slide.shapes[2], "Updated Model Results")
    set_text(
        slide.shapes[3].shapes[1],
        "Forecasting and Clustering\nARIMA forecasting is now implemented, with Denmark's holdout MAPE near 6.57%.\nK-means now identifies three country profiles: hydro-dominant, fast-transition, and low-renewables clusters.",
    )
    set_text(
        slide.shapes[4].shapes[1],
        "Decomposition and Lags\nWind explains most of Denmark and Germany's renewable-share increase, while Japan looks more solar-driven.\nLagged generation models strengthen the timing story, but the joint share model is still weak.",
    )

    # Slide 11
    slide = prs.slides[10]
    set_text(slide.shapes[2], "Lagged Solar Generation Effects")
    replace_picture(slide, slide.shapes[3], LAG_SOLAR_IMG)

    # Slide 12
    slide = prs.slides[11]
    set_text(slide.shapes[2], "Challenges")
    set_text(
        slide.shapes[3].shapes[1],
        "Data Constraints\nThe deepest modeling still depends on only 9 countries with overlapping capacity data, so generalizability remains limited.",
    )
    set_text(
        slide.shapes[4].shapes[1],
        "Interpretation Limits\nLagged models improve the timing story, but they still show associations rather than a fully identified causal effect.",
    )
    set_text(
        slide.shapes[5].shapes[1],
        "Modeling Gaps\nForecasting is implemented for one renewable-share series, but broader validation across countries and targets is still needed.",
    )

    # Slide 13
    slide = prs.slides[12]
    set_text(slide.shapes[2], "Next Steps")
    set_text(
        slide.shapes[3].shapes[1],
        "Expand forecasting\nRun ARIMA across more countries and compare against simple baseline forecasts.",
    )
    set_text(
        slide.shapes[4].shapes[1],
        "Strengthen identification\nAdd better controls or external policy variables if we want a stronger causal interpretation.",
    )
    set_text(
        slide.shapes[5].shapes[1],
        "Refine transition story\nUse decomposition and clustering together to classify hydro-led, wind-led, and solar-led transition paths.",
    )
    set_text(
        slide.shapes[6].shapes[1],
        "Finalize presentation\nSelect the clearest visuals and turn the current outputs into a concise final narrative.",
    )

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
